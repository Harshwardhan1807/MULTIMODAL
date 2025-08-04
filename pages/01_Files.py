import streamlit as st
from langchain_text_splitters import RecursiveCharacterTextSplitter
from PyPDF2 import PdfReader
from langchain_google_genai import GoogleGenerativeAIEmbeddings, ChatGoogleGenerativeAI
from PyPDF2 import PdfReader
from docx import Document as DocxReader
from pptx import Presentation
from dotenv import load_dotenv
from langchain_pinecone import PineconeVectorStore
import os
from langchain import hub
from langchain.chains.combine_documents import create_stuff_documents_chain
from langchain.chains.retrieval import create_retrieval_chain
from pinecone import Pinecone
load_dotenv()
os.environ["ALLOW_RESET"] = "TRUE"


def get_text(uploaded_files):
    text = ""
    for file in uploaded_files:
        file_name = file.name.lower()

        if file_name.endswith(".pdf"):
            reader = PdfReader(file)
            for page in reader.pages:
                text+= page.extract_text() + "\n"
            
            if text.strip() == "" or len(text) < 10:
                text = "No text found in the PDF"

        elif file_name.endswith(".docx"):
            doc = DocxReader(file)
            for para in doc.paragraphs:
                text += para.text + "\n"

        elif file_name.endswith(".pptx"):
            prs = Presentation(file)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"

        elif file_name.endswith(".txt"):
            text += file.read().decode("utf-8") + "\n"
        
    return text

def get_text_chunks(text):
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=3000, chunk_overlap=500)
    chunks = text_splitter.split_text(text)
    return chunks

def get_vector_store(chunks):
    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001", google_api_key=os.getenv("GEMINI_API_KEY")) 
    PineconeVectorStore.from_texts(chunks, embeddings, index_name=os.getenv("PINECONE_INDEX_NAME"))
    vstore = PineconeVectorStore(embedding=embeddings, index_name=os.getenv("PINECONE_INDEX_NAME"))
    return vstore

def user_input(user_question, vstore):
    llm = ChatGoogleGenerativeAI(model="gemini-2.0-flash",temperature=0.3,google_api_key=os.getenv("GEMINI_API_KEY"))
    retrieval_chat_prompt = hub.pull("langchain-ai/retrieval-qa-chat")
    combine_docs_chain = create_stuff_documents_chain(llm=llm, prompt=retrieval_chat_prompt)
    retrieval_chain = create_retrieval_chain(vstore.as_retriever(), combine_docs_chain)
    response = retrieval_chain.invoke({"input": user_question})
    st.write(response["answer"])

def main():
    st.set_page_config("Chat with documents")
    with st.sidebar:
        docs = st.file_uploader("Upload your files", type=["pdf", "docx", "pptx", "txt"],accept_multiple_files=True)
        if st.button("Submit & Process"):
            if docs == []:
                st.error("Please upload files first")
                return
            with st.spinner("Processing..."):
                text = get_text(docs)
                if "No text found in the PDF" in text:
                    st.error("No text found in the PDF")
                    return
                text_chunks = get_text_chunks(text)
                vstore = get_vector_store(text_chunks)
                st.session_state.vstore = vstore
                st.success("Done")
        if st.button("Erase all memory"):
            pc = Pinecone(api_key=os.getenv("PINECONE_API_KEY")) 
            index = pc.Index(os.getenv("PINECONE_INDEX_NAME"))
            index.delete(delete_all=True)
            st.session_state.vstore = None

    st.header("Chat with documents")
    user_question = st.text_input("Ask a Question from the Files")
    if st.button("Submit"):
        if user_question:
            if "vstore" not in st.session_state or st.session_state.vstore is None:
                st.error("Please upload files first")
                return
            with st.spinner(text="In progress..."):
                user_input(user_question, st.session_state.vstore)

    
if __name__ == "__main__": 
    main()